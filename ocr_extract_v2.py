"""
OCR-скрипт v2 — с промежуточным сохранением и обработкой ошибок.
"""
import sys
from pathlib import Path
from io import BytesIO

import fitz
import pytesseract
from PIL import Image

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

BASE_DIR = Path(__file__).parent / "база данных"
OUT_DIR = Path(__file__).parent / "extracted"
OUT_DIR.mkdir(exist_ok=True)


def ocr_pdf(filepath: Path, out_name: str, dpi: int = 250, start_page: int = 0):
    out_path = OUT_DIR / out_name
    doc = fitz.open(str(filepath))
    total = len(doc)
    print(f"OCR: {filepath.name} — {total} pages, DPI={dpi}, starting from page {start_page + 1}")

    # Load existing text if resuming
    existing = ""
    if start_page > 0 and out_path.exists():
        existing = out_path.read_text(encoding="utf-8")

    for i in range(start_page, total):
        try:
            page = doc[i]
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            img = Image.open(BytesIO(pix.tobytes("png")))
            text = pytesseract.image_to_string(img, lang="rus+eng", config="--psm 6")
            chars = len(text.strip())
            print(f"  [{i+1}/{total}] {chars} chars")
            if text.strip():
                existing += f"\n\n--- Страница {i+1} ---\n{text.strip()}"
        except (KeyboardInterrupt, SystemExit):
            print(f"\n  [INTERRUPTED at page {i+1}] Saving progress...")
            out_path.write_text(existing, encoding="utf-8")
            print(f"  [SAVED] {len(existing)} chars -> {out_path}")
            doc.close()
            sys.exit(1)
        except Exception as e:
            print(f"  [{i+1}/{total}] ERROR: {e}")
            existing += f"\n\n--- Страница {i+1} ---\n[ОШИБКА OCR: {e}]"

        # Save every 10 pages
        if (i + 1) % 10 == 0 or i == total - 1:
            out_path.write_text(existing, encoding="utf-8")
            print(f"  [SAVED] {len(existing)} chars -> {out_path}")

    doc.close()
    print(f"\n[DONE] {out_path} — {len(existing)} chars")


def ocr_pptx(filepath: Path, out_name: str):
    from pptx import Presentation
    
    out_path = OUT_DIR / out_name
    prs = Presentation(str(filepath))
    all_text = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.shape_type == 13:
                try:
                    img = Image.open(BytesIO(shape.image.blob))
                    if img.width > 100 and img.height > 100:
                        text = pytesseract.image_to_string(img, lang="rus+eng", config="--psm 6")
                        if text.strip():
                            slide_texts.append(text.strip())
                except Exception as e:
                    print(f"  Slide {slide_num}: img error: {e}")
            if hasattr(shape, 'shapes'):
                for sub in shape.shapes:
                    if hasattr(sub, 'text_frame'):
                        for para in sub.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
        
        print(f"  Slide {slide_num}: {sum(len(t) for t in slide_texts)} chars")
        if slide_texts:
            all_text.append(f"\n--- Слайд {slide_num} ---\n" + "\n".join(slide_texts))
    
    result = "\n".join(all_text)
    out_path.write_text(result, encoding="utf-8")
    print(f"[DONE] {out_path} — {len(result)} chars")


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--start", type=int, default=0, help="Start page (0-based)")
    parser.add_argument("--skip-pdf", action="store_true")
    parser.add_argument("--skip-pptx", action="store_true")
    args = parser.parse_args()

    if not args.skip_pdf:
        pdf1 = BASE_DIR / "Справочник_Fitline_2024_2025.pdf"
        if pdf1.exists():
            ocr_pdf(pdf1, "spravochnik_ocr.txt", dpi=250, start_page=args.start)

    if not args.skip_pptx:
        pptx1 = BASE_DIR / "ПРОДУКТ 2025.pptx"
        if pptx1.exists():
            ocr_pptx(pptx1, "product_2025_ocr.txt")
