"""
Альтернативный скрипт: пробуем разные методы извлечения текста.
- PyPDF2 для PDF
- python-pptx с проверкой notes, grouped shapes и images
- pypdfium2 для рендера страниц (если OCR нужен)
"""
from pathlib import Path
import sys

BASE_DIR = Path(__file__).parent / "база данных"
OUT_DIR = Path(__file__).parent / "extracted"
OUT_DIR.mkdir(exist_ok=True)


def try_pypdf2(filepath: Path, out_name: str):
    """Try extracting with PyPDF2"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        print("[PyPDF2] not installed, skipping")
        return
    reader = PdfReader(str(filepath))
    text_parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            text_parts.append(f"\n--- Страница {i} ---\n{text.strip()}")
    result = "\n".join(text_parts)
    if result.strip():
        out_path = OUT_DIR / out_name
        out_path.write_text(result, encoding="utf-8")
        print(f"[PyPDF2 OK] {filepath.name} -> {out_path} ({len(result)} chars)")
    else:
        print(f"[PyPDF2] {filepath.name}: 0 chars extracted (likely scanned/image-based)")
    return result


def try_pypdfium2(filepath: Path, out_name: str):
    """Try extracting with pypdfium2"""
    try:
        import pypdfium2 as pdfium
    except ImportError:
        print("[pypdfium2] not installed, skipping")
        return
    pdf = pdfium.PdfDocument(str(filepath))
    text_parts = []
    for i in range(len(pdf)):
        page = pdf[i]
        textpage = page.get_textpage()
        text = textpage.get_text_range()
        if text and text.strip():
            text_parts.append(f"\n--- Страница {i+1} ---\n{text.strip()}")
        textpage.close()
        page.close()
    pdf.close()
    result = "\n".join(text_parts)
    if result.strip():
        out_path = OUT_DIR / out_name
        out_path.write_text(result, encoding="utf-8")
        print(f"[pypdfium2 OK] {filepath.name} -> {out_path} ({len(result)} chars)")
    else:
        print(f"[pypdfium2] {filepath.name}: 0 chars (image-based PDF)")
    return result


def try_pptx_deep(filepath: Path, out_name: str):
    """Deep extraction from PPTX: shapes, groups, notes, tables"""
    from pptx import Presentation
    from pptx.util import Inches
    import pptx.oxml.ns as ns

    prs = Presentation(str(filepath))
    text_parts = []

    print(f"[PPTX] {filepath.name}: {len(prs.slides)} slides")

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        # Count shapes
        shape_count = len(slide.shapes)

        for shape in slide.shapes:
            shape_type = type(shape).__name__
            # Text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(cells))
            # Group shapes
            if hasattr(shape, 'shapes'):
                for sub in shape.shapes:
                    if hasattr(sub, 'text_frame'):
                        for para in sub.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)

        # Notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            for para in notes.paragraphs:
                t = para.text.strip()
                if t:
                    slide_texts.append(f"[Заметка] {t}")

        if slide_texts:
            text_parts.append(f"\n--- Слайд {slide_num} ({shape_count} shapes) ---\n" + "\n".join(slide_texts))
        else:
            text_parts.append(f"\n--- Слайд {slide_num} ({shape_count} shapes) --- [только изображения]")

    result = "\n".join(text_parts)
    out_path = OUT_DIR / out_name
    out_path.write_text(result, encoding="utf-8")
    print(f"[PPTX OK] {filepath.name} -> {out_path} ({len(result)} chars)")
    return result


if __name__ == "__main__":
    pdf1 = BASE_DIR / "Справочник_Fitline_2024_2025.pdf"
    pptx1 = BASE_DIR / "ПРОДУКТ 2025.pptx"

    print("=== Попытка 2: альтернативные методы ===\n")

    if pdf1.exists():
        print(f"--- {pdf1.name} ---")
        try_pypdf2(pdf1, "spravochnik_pypdf2.txt")
        try_pypdfium2(pdf1, "spravochnik_pypdfium2.txt")
    
    if pptx1.exists():
        print(f"\n--- {pptx1.name} ---")
        try_pptx_deep(pptx1, "product_2025_deep.txt")

    print("\n=== Готово ===")
