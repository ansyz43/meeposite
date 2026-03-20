"""
OCR-скрипт для извлечения текста из отсканированного PDF и PPTX с изображениями.
Использует PyMuPDF (fitz) для рендера PDF → изображение, pytesseract + Tesseract OCR для распознавания.
"""
import sys
import os
from pathlib import Path
from io import BytesIO

import fitz  # PyMuPDF
import pytesseract
from PIL import Image

# Путь к Tesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

BASE_DIR = Path(__file__).parent / "база данных"
OUT_DIR = Path(__file__).parent / "extracted"
OUT_DIR.mkdir(exist_ok=True)


def ocr_pdf(filepath: Path, out_name: str, dpi: int = 300):
    """Render each PDF page to image and OCR with Tesseract (Russian + English)."""
    print(f"\n{'='*60}")
    print(f"OCR PDF: {filepath.name}")
    print(f"{'='*60}")

    doc = fitz.open(str(filepath))
    total_pages = len(doc)
    print(f"Pages: {total_pages}")

    all_text = []
    for page_num in range(total_pages):
        page = doc[page_num]
        # Render page to image at specified DPI
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)

        # Convert to PIL Image
        img_data = pix.tobytes("png")
        img = Image.open(BytesIO(img_data))

        # OCR with Russian + English
        text = pytesseract.image_to_string(img, lang="rus+eng", config="--psm 6")

        progress = f"[{page_num + 1}/{total_pages}]"
        chars = len(text.strip())
        print(f"  {progress} Page {page_num + 1}: {chars} chars")

        if text.strip():
            all_text.append(f"\n--- Страница {page_num + 1} ---\n{text.strip()}")

    doc.close()

    result = "\n".join(all_text)
    out_path = OUT_DIR / out_name
    out_path.write_text(result, encoding="utf-8")
    print(f"\n[DONE] {out_path} — {len(result)} chars total")
    return result


def ocr_pptx(filepath: Path, out_name: str):
    """Extract images from PPTX and OCR them."""
    from pptx import Presentation

    print(f"\n{'='*60}")
    print(f"OCR PPTX: {filepath.name}")
    print(f"{'='*60}")

    prs = Presentation(str(filepath))
    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        for shape in slide.shapes:
            # Text from text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)

            # OCR images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_blob = shape.image.blob
                    img = Image.open(BytesIO(img_blob))
                    # Only OCR if image is large enough to contain text
                    if img.width > 100 and img.height > 100:
                        text = pytesseract.image_to_string(img, lang="rus+eng", config="--psm 6")
                        if text.strip():
                            slide_texts.append(f"[Изображение] {text.strip()}")
                except Exception as e:
                    print(f"  Slide {slide_num}: image error: {e}")

            # Group shapes
            if hasattr(shape, 'shapes'):
                for sub in shape.shapes:
                    if hasattr(sub, 'text_frame'):
                        for para in sub.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                    if hasattr(sub, 'shape_type') and sub.shape_type == 13:
                        try:
                            img_blob = sub.image.blob
                            img = Image.open(BytesIO(img_blob))
                            if img.width > 100 and img.height > 100:
                                text = pytesseract.image_to_string(img, lang="rus+eng", config="--psm 6")
                                if text.strip():
                                    slide_texts.append(f"[Изображение] {text.strip()}")
                        except Exception:
                            pass

        chars = sum(len(t) for t in slide_texts)
        print(f"  Slide {slide_num}: {chars} chars, {len(slide_texts)} blocks")

        if slide_texts:
            all_text.append(f"\n--- Слайд {slide_num} ---\n" + "\n".join(slide_texts))

    result = "\n".join(all_text)
    out_path = OUT_DIR / out_name
    out_path.write_text(result, encoding="utf-8")
    print(f"\n[DONE] {out_path} — {len(result)} chars total")
    return result


if __name__ == "__main__":
    print("=== OCR Extraction (Tesseract + PyMuPDF) ===")
    print(f"Tesseract: {pytesseract.pytesseract.tesseract_cmd}")

    # PDF - Справочник
    pdf1 = BASE_DIR / "Справочник_Fitline_2024_2025.pdf"
    if pdf1.exists():
        ocr_pdf(pdf1, "spravochnik_ocr.txt", dpi=300)
    else:
        print(f"[SKIP] {pdf1}")

    # PPTX - Продукт 2025
    pptx1 = BASE_DIR / "ПРОДУКТ 2025.pptx"
    if pptx1.exists():
        ocr_pptx(pptx1, "product_2025_ocr.txt")
    else:
        print(f"[SKIP] {pptx1}")

    print("\n=== Готово! ===")
