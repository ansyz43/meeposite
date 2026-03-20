"""
Одноразовый скрипт для извлечения текста из PDF и PPTX файлов
из папки 'база данных' для обновления базы знаний бота.
"""
import sys
from pathlib import Path

BASE_DIR = Path(__file__).parent / "база данных"
OUT_DIR = Path(__file__).parent / "extracted"
OUT_DIR.mkdir(exist_ok=True)


def extract_pdf(filepath: Path, out_name: str):
    import pdfplumber
    text_parts = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"\n--- Страница {i} ---\n{page_text}")
            # Also try extracting tables
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    if row:
                        cleaned = [str(cell).strip() if cell else "" for cell in row]
                        text_parts.append(" | ".join(cleaned))
    result = "\n".join(text_parts)
    out_path = OUT_DIR / out_name
    out_path.write_text(result, encoding="utf-8")
    print(f"[OK] {filepath.name} -> {out_path} ({len(result)} chars, {len(text_parts)} blocks)")
    return result


def extract_pptx(filepath: Path, out_name: str):
    from pptx import Presentation
    text_parts = []
    prs = Presentation(str(filepath))
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(cells))
        if slide_texts:
            text_parts.append(f"\n--- Слайд {slide_num} ---\n" + "\n".join(slide_texts))
    result = "\n".join(text_parts)
    out_path = OUT_DIR / out_name
    out_path.write_text(result, encoding="utf-8")
    print(f"[OK] {filepath.name} -> {out_path} ({len(result)} chars, {len(text_parts)} slides)")
    return result


if __name__ == "__main__":
    print("=== Извлечение текста из PDF/PPTX ===\n")

    pdf1 = BASE_DIR / "Справочник_Fitline_2024_2025.pdf"
    pdf2 = BASE_DIR / "Спортивный каталог.pdf"
    pptx1 = BASE_DIR / "ПРОДУКТ 2025.pptx"

    if pdf1.exists():
        extract_pdf(pdf1, "spravochnik.txt")
    else:
        print(f"[SKIP] {pdf1} not found")

    if pdf2.exists():
        extract_pdf(pdf2, "sport_catalog.txt")
    else:
        print(f"[SKIP] {pdf2} not found")

    if pptx1.exists():
        extract_pptx(pptx1, "product_2025.txt")
    else:
        print(f"[SKIP] {pptx1} not found")

    print("\n=== Готово! Файлы в папке 'extracted/' ===")
